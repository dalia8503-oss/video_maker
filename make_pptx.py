"""
AI 영상 제작 시스템 소개 PPTX 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 색상 팔레트 ──────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # #0d1117  dark background
ACCENT_CYAN  = RGBColor(0x00, 0xB4, 0xD8)   # #00b4d8  titles / highlights
ACCENT_BLUE  = RGBColor(0x4F, 0xC3, 0xF7)   # #4fc3f7  sub-highlights
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)   # body text
CARD_BG      = RGBColor(0x16, 0x1B, 0x22)   # #161b22  card background
CARD_BORDER  = RGBColor(0x21, 0x26, 0x2D)   # subtle border
GREEN_ACCENT = RGBColor(0x3F, 0xB9, 0x50)   # #3fb950
ORANGE_ACC   = RGBColor(0xF7, 0x8A, 0x1B)   # #f78a1b

# ── 슬라이드 크기 (16:9 standard) ────────────────────────────
SLIDE_W = Emu(9144000)   # 13.33"
SLIDE_H = Emu(5143500)   # 7.5"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # blank

# ════════════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════════════

def add_bg(slide, color=BG_DARK):
    """슬라이드 전체 배경 사각형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def txb(slide, text, left, top, width, height,
        font_size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    """텍스트 박스 헬퍼"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"
    return box


def txb_lines(slide, lines, left, top, width, height,
              font_size=18, bold=False, color=WHITE,
              line_spacing_pt=None, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트 박스 (각 줄이 별도 paragraph)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing_pt * 100)))
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Malgun Gothic"
    return box


def add_accent_line(slide, top, color=ACCENT_CYAN, thickness=Pt(3)):
    """수평 강조선"""
    from pptx.util import Emu
    line = slide.shapes.add_shape(1,
        Emu(int(SLIDE_W * 0.05)), top,
        Emu(int(SLIDE_W * 0.90)), thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def slide_title(slide, title_text, top=Emu(220000)):
    """공통 슬라이드 제목"""
    txb(slide, title_text,
        left=Emu(int(SLIDE_W * 0.05)),
        top=top,
        width=Emu(int(SLIDE_W * 0.90)),
        height=Emu(500000),
        font_size=36, bold=True, color=ACCENT_CYAN,
        align=PP_ALIGN.LEFT)
    add_accent_line(slide, top + Emu(530000))


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — 요약 (Executive Summary)
# ════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide1)

# ── 상단 타이틀 영역 ──
add_rect(slide1,
    Emu(0), Emu(0),
    SLIDE_W, Emu(1400000),
    RGBColor(0x0A, 0x0E, 0x1A))

txb(slide1, "🎬 AI 영상 제작 시스템",
    Emu(int(SLIDE_W*0.05)), Emu(120000),
    Emu(int(SLIDE_W*0.90)), Emu(600000),
    font_size=44, bold=True, color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER)

txb(slide1, "사진 한 장이 추억 영상이 되는 순간",
    Emu(int(SLIDE_W*0.05)), Emu(720000),
    Emu(int(SLIDE_W*0.90)), Emu(400000),
    font_size=22, bold=False, color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER)

# ── 5개 키포인트 카드 ──
cards = [
    ("📸", "사진 → 영상\n자동 변환"),
    ("🤖", "GPT-4o Vision\n자막 생성"),
    ("✂️", "음악 구간\n편집"),
    ("✨", "AI 자막\n감성 다듬기"),
    ("🎞️", "HD MP4\n즉시 다운로드"),
]

n = len(cards)
card_w = Emu(int(SLIDE_W * 0.16))
card_h = Emu(int(SLIDE_H * 0.38))
gap = Emu(int((SLIDE_W * 0.90 - card_w * n) / (n - 1)))
left_start = Emu(int(SLIDE_W * 0.05))
card_top = Emu(1530000)

for i, (icon, label) in enumerate(cards):
    cx = left_start + i * (card_w + gap)
    # card bg
    card = add_rect(slide1, cx, card_top, card_w, card_h,
                    CARD_BG, ACCENT_CYAN, 1.5)
    # icon
    txb(slide1, icon,
        cx, card_top + Emu(100000), card_w, Emu(550000),
        font_size=38, align=PP_ALIGN.CENTER)
    # label
    txb(slide1, label,
        cx, card_top + Emu(600000), card_w, Emu(500000),
        font_size=15, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)

# ── 기술 스택 한 줄 ──
txb(slide1,
    "Python  ·  Streamlit  ·  OpenAI GPT-4o Vision  ·  MoviePy + FFmpeg  ·  PIL/Pillow",
    Emu(int(SLIDE_W*0.05)), Emu(3600000),
    Emu(int(SLIDE_W*0.90)), Emu(320000),
    font_size=15, color=ACCENT_BLUE,
    align=PP_ALIGN.CENTER)

# ── 하단 'made by' ──
add_rect(slide1,
    Emu(0), Emu(int(SLIDE_H - 340000)),
    SLIDE_W, Emu(340000),
    RGBColor(0x08, 0x0C, 0x14))
txb(slide1, "made by  s.y.Kim",
    Emu(int(SLIDE_W*0.05)), Emu(int(SLIDE_H - 310000)),
    Emu(int(SLIDE_W*0.90)), Emu(280000),
    font_size=14, color=LIGHT_GRAY,
    align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — 시스템 구성
# ════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide2)
slide_title(slide2, "🏗️ 시스템 구성")

# 아키텍처 플로우 박스
flow_items = [
    ("Streamlit UI", "웹 인터페이스\n(브라우저)", ACCENT_BLUE),
    ("Python Backend", "비즈니스 로직\n상태 관리", ACCENT_CYAN),
    ("OpenAI API", "GPT-4o Vision\n자막 생성", GREEN_ACCENT),
    ("FFmpeg / MoviePy", "영상·음성\n합성 처리", ORANGE_ACC),
]

fw = Emu(int(SLIDE_W * 0.18))
fh = Emu(int(SLIDE_H * 0.22))
fgap = Emu(int((SLIDE_W * 0.85 - fw * 4) / 3))
fl_start = Emu(int(SLIDE_W * 0.075))
ft = Emu(1000000)

arrow_top = ft + Emu(int(fh / 2)) - Emu(60000)

for i, (title, desc, color) in enumerate(flow_items):
    fx = fl_start + i * (fw + fgap)
    add_rect(slide2, fx, ft, fw, fh, CARD_BG, color, 2)
    txb(slide2, title, fx, ft + Emu(60000), fw, Emu(250000),
        font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide2, desc, fx, ft + Emu(320000), fw, Emu(350000),
        font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(flow_items) - 1:
        ax = fx + fw + Emu(30000)
        txb(slide2, "→", ax, arrow_top, fgap - Emu(60000), Emu(220000),
            font_size=22, bold=True, color=ACCENT_CYAN,
            align=PP_ALIGN.CENTER)

# 컴포넌트 테이블
comp_top = Emu(2380000)
headers = ["구성 요소", "기술", "역할"]
col_w   = [Emu(int(SLIDE_W*0.22)), Emu(int(SLIDE_W*0.30)), Emu(int(SLIDE_W*0.38))]
col_x   = [Emu(int(SLIDE_W*0.05))]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + Emu(20000))
row_h = Emu(310000)

# header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide2, cx, comp_top, cw, row_h,
             RGBColor(0x16, 0x35, 0x50), ACCENT_CYAN, 1)
    txb(slide2, hdr, cx + Emu(30000), comp_top + Emu(60000),
        cw - Emu(60000), row_h - Emu(60000),
        font_size=15, bold=True, color=ACCENT_CYAN)

rows = [
    ("Frontend",          "Streamlit",             "웹 인터페이스 제공"),
    ("AI Engine",         "OpenAI GPT-4o Vision",  "이미지 분석·자막 생성"),
    ("Video Engine",      "MoviePy + FFmpeg",       "영상 합성·트랜지션"),
    ("Image Processing",  "PIL / Pillow",           "자막 렌더링·이모지 합성"),
    ("Audio",             "FFmpeg 직접 호출",        "음악 트리밍·루프·페이드아웃"),
]

for r, (c0, c1, c2) in enumerate(rows):
    ry = comp_top + (r + 1) * row_h
    row_bg = CARD_BG if r % 2 == 0 else RGBColor(0x11, 0x16, 0x1D)
    for j, (cell, cx, cw) in enumerate(zip([c0, c1, c2], col_x, col_w)):
        add_rect(slide2, cx, ry, cw, row_h, row_bg,
                 RGBColor(0x21, 0x26, 0x2D), 0.5)
        cell_color = ACCENT_BLUE if j == 1 else LIGHT_GRAY
        txb(slide2, cell, cx + Emu(30000), ry + Emu(60000),
            cw - Emu(60000), row_h - Emu(60000),
            font_size=14, color=cell_color)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — 사진 관리 기능
# ════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide3)
slide_title(slide3, "📸 사진 관리 기능")

features = [
    ("🔼🔽", "순서 변경",   "화살표 버튼으로 사진 순서를 원하는 대로 자유롭게 조정"),
    ("🔄",   "개별 회전",   "사진마다 90° 독립 회전 — 촬영 방향 오류 즉시 교정"),
    ("🗑️",   "개별 삭제",   "미리보기에서 불필요한 사진을 클릭 한 번으로 제거"),
    ("📁",   "폴더 불러오기","서버 폴더 경로 입력만으로 대량 사진 일괄 로드"),
    ("👁️",   "실시간 미리보기","모든 변경사항이 즉시 반영되어 결과 확인 가능"),
]

feat_top = Emu(1050000)
feat_h   = Emu(640000)
feat_gap = Emu(60000)
icon_w   = Emu(int(SLIDE_W * 0.08))
label_w  = Emu(int(SLIDE_W * 0.16))
desc_w   = Emu(int(SLIDE_W * 0.62))
left0    = Emu(int(SLIDE_W * 0.05))

for i, (icon, label, desc) in enumerate(features):
    fy = feat_top + i * (feat_h + feat_gap)
    # row bg
    add_rect(slide3, left0, fy,
             icon_w + label_w + desc_w + Emu(40000), feat_h,
             CARD_BG, RGBColor(0x21, 0x26, 0x2D), 0.8)
    # icon
    txb(slide3, icon, left0, fy + Emu(80000), icon_w, feat_h - Emu(80000),
        font_size=30, align=PP_ALIGN.CENTER)
    # label
    txb(slide3, label,
        left0 + icon_w, fy + Emu(110000),
        label_w, feat_h - Emu(110000),
        font_size=17, bold=True, color=ACCENT_CYAN)
    # desc
    txb(slide3, desc,
        left0 + icon_w + label_w, fy + Emu(110000),
        desc_w, feat_h - Emu(110000),
        font_size=15, color=LIGHT_GRAY)

# 우측 업로드 방식 카드
card_right_x = Emu(int(SLIDE_W * 0.77))
card_right_w = Emu(int(SLIDE_W * 0.18))
add_rect(slide3,
    card_right_x, feat_top,
    card_right_w, Emu(int((feat_h + feat_gap) * 5 - feat_gap)),
    RGBColor(0x10, 0x20, 0x30), ACCENT_BLUE, 1.5)

txb_lines(slide3,
    ["업로드 방식", "",
     "① 브라우저\n   직접 업로드",
     "",
     "② 서버 폴더\n   경로 지정"],
    card_right_x + Emu(40000), feat_top + Emu(80000),
    card_right_w - Emu(80000),
    Emu(int((feat_h + feat_gap) * 5 - feat_gap) - Emu(100000)),
    font_size=13, color=LIGHT_GRAY)

# 제목 강조 재설정 (오버라이드)
txb(slide3, "업로드 방식",
    card_right_x + Emu(40000), feat_top + Emu(80000),
    card_right_w - Emu(80000), Emu(280000),
    font_size=14, bold=True, color=ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — AI 자막 생성 (핵심 슬라이드)
# ════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide4)
slide_title(slide4, "🤖 AI 자막 생성 — GPT-4o Vision API")

# API 하이라이트 배너
add_rect(slide4,
    Emu(int(SLIDE_W*0.05)), Emu(820000),
    Emu(int(SLIDE_W*0.90)), Emu(280000),
    RGBColor(0x00, 0x30, 0x50), ACCENT_CYAN, 1.5)
txb(slide4,
    "📡  사진을 직접 분석하여 맥락에 맞는 자막을 자동 생성  |  Token 사용량 실시간 추적",
    Emu(int(SLIDE_W*0.06)), Emu(840000),
    Emu(int(SLIDE_W*0.88)), Emu(240000),
    font_size=15, bold=True, color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER)

# 4개 스타일 카드
style_cards = [
    ("인스타그램 인플루언서",
     "완벽한 여름 한 컷 ☀️",
     ACCENT_CYAN),
    ("카피라이터 / 유머",
     "아직 여기 있고 싶은데\n집이 날 부른다",
     GREEN_ACCENT),
    ("감성 브이로그",
     "완벽했던 오후",
     ACCENT_BLUE),
    ("숏폼 여행 인플루언서",
     "윤슬 미쳤다...\n여기가 바로 무릉도원 🌊",
     ORANGE_ACC),
]

sc_top  = Emu(1230000)
sc_w    = Emu(int(SLIDE_W * 0.205))
sc_h    = Emu(int(SLIDE_H * 0.44))
sc_gap  = Emu(int((SLIDE_W * 0.90 - sc_w * 4) / 3))
sc_left = Emu(int(SLIDE_W * 0.05))

for i, (style_name, example, color) in enumerate(style_cards):
    sx = sc_left + i * (sc_w + sc_gap)
    add_rect(slide4, sx, sc_top, sc_w, sc_h, CARD_BG, color, 2)
    # style number
    txb(slide4, f"Style {i+1}",
        sx, sc_top + Emu(50000), sc_w, Emu(240000),
        font_size=12, bold=False, color=color,
        align=PP_ALIGN.CENTER)
    # style name
    txb(slide4, style_name,
        sx, sc_top + Emu(220000), sc_w, Emu(300000),
        font_size=14, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    # divider
    add_rect(slide4,
        sx + Emu(int(sc_w*0.10)), sc_top + Emu(520000),
        Emu(int(sc_w*0.80)), Emu(30000), color)
    # example
    txb(slide4, f'"{example}"',
        sx + Emu(40000), sc_top + Emu(590000),
        sc_w - Emu(80000), sc_h - Emu(650000),
        font_size=13, bold=False, color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER)

# 하단 토큰 사용량 표시
token_top = Emu(3700000)
add_rect(slide4,
    Emu(int(SLIDE_W*0.05)), token_top,
    Emu(int(SLIDE_W*0.90)), Emu(350000),
    RGBColor(0x14, 0x20, 0x10), GREEN_ACCENT, 1)
txb_lines(slide4,
    ["📊  Token 사용량 추적",
     "자막 생성 시 사용된 Prompt / Completion / Total Token을 실시간으로 표시 — 비용 관리 가능"],
    Emu(int(SLIDE_W*0.07)), token_top + Emu(40000),
    Emu(int(SLIDE_W*0.86)), Emu(290000),
    font_size=14, color=LIGHT_GRAY)

txb(slide4, "📊  Token 사용량 추적",
    Emu(int(SLIDE_W*0.07)), token_top + Emu(40000),
    Emu(int(SLIDE_W*0.30)), Emu(210000),
    font_size=14, bold=True, color=GREEN_ACCENT)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — 자막 편집 & AI 다듬기
# ════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide5)
slide_title(slide5, "✍️ 자막 편집 & AI 다듬기")

# 2단계 워크플로우 다이어그램
wf_top = Emu(1020000)
wf_box_w = Emu(int(SLIDE_W * 0.25))
wf_box_h = Emu(600000)
wf_gap   = Emu(int(SLIDE_W * 0.06))
wf_left  = Emu(int(SLIDE_W * 0.06))

stages = [
    ("Stage 1", "🤖 AI 자막\n자동 생성", ACCENT_CYAN),
    ("Stage 2", "✏️ 사용자\n직접 수정", ACCENT_BLUE),
    ("✨ AI 다듬기", "감성 개선 +\n이모지 삽입", ORANGE_ACC),
    ("🎬 최종 영상", "확정된 자막\n으로 렌더링", GREEN_ACCENT),
]

for i, (stg, desc, color) in enumerate(stages):
    sx = wf_left + i * (wf_box_w + wf_gap)
    add_rect(slide5, sx, wf_top, wf_box_w, wf_box_h,
             CARD_BG, color, 2)
    txb(slide5, stg, sx, wf_top + Emu(50000), wf_box_w, Emu(250000),
        font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide5, desc, sx, wf_top + Emu(280000), wf_box_w, Emu(280000),
        font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        ax = sx + wf_box_w + Emu(20000)
        txb(slide5, "→", ax, wf_top + Emu(200000),
            wf_gap - Emu(40000), Emu(250000),
            font_size=24, bold=True, color=ACCENT_CYAN,
            align=PP_ALIGN.CENTER)

# AI 다듬기 기능 세부 설명
detail_top = Emu(1900000)
add_rect(slide5,
    Emu(int(SLIDE_W*0.05)), detail_top,
    Emu(int(SLIDE_W*0.55)), Emu(1200000),
    CARD_BG, ORANGE_ACC, 1.5)

txb(slide5, "✨  AI로 자막 다듬기",
    Emu(int(SLIDE_W*0.07)), detail_top + Emu(60000),
    Emu(int(SLIDE_W*0.50)), Emu(240000),
    font_size=18, bold=True, color=ORANGE_ACC)

polish_items = [
    "→  입력한 자막을 감성적으로 개선",
    "→  분위기에 맞는 이모지 자동 삽입",
    "→  원문의 의미·톤을 최대한 유지",
    "→  GPT-4o 프롬프트 엔지니어링 적용",
]
txb_lines(slide5, polish_items,
    Emu(int(SLIDE_W*0.07)), detail_top + Emu(320000),
    Emu(int(SLIDE_W*0.50)), Emu(830000),
    font_size=15, color=LIGHT_GRAY)

# 우측: 수동 입력 모드
add_rect(slide5,
    Emu(int(SLIDE_W*0.62)), detail_top,
    Emu(int(SLIDE_W*0.33)), Emu(1200000),
    CARD_BG, ACCENT_BLUE, 1.5)
txb(slide5, "✏️  수동 입력 모드",
    Emu(int(SLIDE_W*0.64)), detail_top + Emu(60000),
    Emu(int(SLIDE_W*0.29)), Emu(240000),
    font_size=17, bold=True, color=ACCENT_BLUE)
manual_lines = [
    "AI 없이 직접 자막 입력",
    "",
    "• 자유로운 텍스트 입력",
    "• 인터넷 연결 불필요",
    "• API 비용 절감",
]
txb_lines(slide5, manual_lines,
    Emu(int(SLIDE_W*0.64)), detail_top + Emu(320000),
    Emu(int(SLIDE_W*0.29)), Emu(830000),
    font_size=14, color=LIGHT_GRAY)

# 하단 표시 방식
display_top = Emu(3300000)
add_rect(slide5,
    Emu(int(SLIDE_W*0.05)), display_top,
    Emu(int(SLIDE_W*0.90)), Emu(280000),
    RGBColor(0x10, 0x14, 0x24), ACCENT_BLUE, 1)
txb(slide5,
    "🖼️  자막 표시:  영상 하단 반투명 검정 배경  ·  흰색 굵은 텍스트  ·  이모지 포함 Korean Font 렌더링",
    Emu(int(SLIDE_W*0.07)), display_top + Emu(50000),
    Emu(int(SLIDE_W*0.86)), Emu(220000),
    font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — 배경음악 & 영상 출력
# ════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide6)
slide_title(slide6, "🎵 배경음악 & 영상 출력")

# 왼쪽: 배경음악
left_x  = Emu(int(SLIDE_W*0.05))
left_w  = Emu(int(SLIDE_W*0.43))
sect_top = Emu(980000)
sect_h   = Emu(int(SLIDE_H*0.68))

add_rect(slide6, left_x, sect_top, left_w, sect_h, CARD_BG, ACCENT_CYAN, 1.5)
txb(slide6, "🎵  배경음악 기능",
    left_x + Emu(40000), sect_top + Emu(60000),
    left_w - Emu(80000), Emu(260000),
    font_size=20, bold=True, color=ACCENT_CYAN)

music_items = [
    ("🎵", "MP3 / WAV 파일 업로드"),
    ("✂️", "원하는 구간만 선택\n(시작 / 종료 초 단위 설정)"),
    ("🔁", "자동 루프\n(영상 길이에 맞게 반복)"),
    ("🔉", "끝부분 자동 페이드아웃"),
]

mi_top = sect_top + Emu(360000)
for icon, text in music_items:
    add_rect(slide6,
        left_x + Emu(40000), mi_top,
        left_w - Emu(80000), Emu(290000),
        RGBColor(0x0D, 0x20, 0x28), ACCENT_CYAN, 0.8)
    txb(slide6, icon,
        left_x + Emu(60000), mi_top + Emu(50000),
        Emu(300000), Emu(220000),
        font_size=20)
    txb(slide6, text,
        left_x + Emu(380000), mi_top + Emu(50000),
        left_w - Emu(460000), Emu(220000),
        font_size=14, color=LIGHT_GRAY)
    mi_top += Emu(330000)

# 오른쪽: 영상 출력 스펙
right_x = Emu(int(SLIDE_W*0.52))
right_w = Emu(int(SLIDE_W*0.43))
add_rect(slide6, right_x, sect_top, right_w, sect_h, CARD_BG, GREEN_ACCENT, 1.5)
txb(slide6, "🎞️  영상 출력 사양",
    right_x + Emu(40000), sect_top + Emu(60000),
    right_w - Emu(80000), Emu(260000),
    font_size=20, bold=True, color=GREEN_ACCENT)

spec_items = [
    ("📐", "해상도",   "1280 × 720  (HD 720p)"),
    ("🎞️", "포맷",    "MP4  (H.264 + AAC)"),
    ("⚡", "클립 전환", "크로스페이드 트랜지션"),
    ("⬇️", "다운로드", "브라우저 즉시 다운로드"),
]

si_top = sect_top + Emu(360000)
for icon, label, value in spec_items:
    add_rect(slide6,
        right_x + Emu(40000), si_top,
        right_w - Emu(80000), Emu(290000),
        RGBColor(0x0D, 0x22, 0x18), GREEN_ACCENT, 0.8)
    txb(slide6, icon,
        right_x + Emu(60000), si_top + Emu(50000),
        Emu(260000), Emu(220000), font_size=20)
    txb(slide6, label,
        right_x + Emu(340000), si_top + Emu(50000),
        Emu(500000), Emu(220000),
        font_size=13, bold=True, color=GREEN_ACCENT)
    txb(slide6, value,
        right_x + Emu(850000), si_top + Emu(50000),
        right_w - Emu(920000), Emu(220000),
        font_size=14, color=WHITE)
    si_top += Emu(330000)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — 기술적 특장점
# ════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide7)
slide_title(slide7, "🌟 기술적 특장점")

tech_cards = [
    ("🌈", "이모지 렌더링",
     "Twemoji PNG 다운로드 →\n투명 패딩 제거 → 정밀 합성",
     ORANGE_ACC),
    ("🔐", "비밀번호 인증",
     "시스템 접근 보안\n(세션 기반 인증 관리)",
     RGBColor(0xF8, 0x5B, 0x53)),
    ("📦", "자동 패키지 설치",
     "실행 시 의존성 자동 처리\n(subprocess pip install)",
     GREEN_ACCENT),
    ("🗂️", "2단계 파이프라인",
     "자막 편집 → 최종 렌더링\n명확한 단계 분리",
     ACCENT_CYAN),
    ("💾", "세션 상태 관리",
     "Streamlit session_state 기반\n상태 영속성 유지",
     ACCENT_BLUE),
]

tc_per_row = 3
tc_w = Emu(int(SLIDE_W * 0.27))
tc_h = Emu(int(SLIDE_H * 0.28))
tc_h_gap = Emu(int(SLIDE_W * 0.025))
tc_v_gap = Emu(100000)
tc_top  = Emu(1000000)
tc_left = Emu(int((SLIDE_W - tc_w * tc_per_row - tc_h_gap * (tc_per_row - 1)) / 2))

for i, (icon, title, desc, color) in enumerate(tech_cards):
    row = i // tc_per_row
    col = i % tc_per_row
    if i == 3:
        # 2nd row: center 2 cards
        tc_left2 = Emu(int((SLIDE_W - tc_w * 2 - tc_h_gap) / 2))
        tx = tc_left2 + col * (tc_w + tc_h_gap)
    elif i == 4:
        tc_left2 = Emu(int((SLIDE_W - tc_w * 2 - tc_h_gap) / 2))
        tx = tc_left2 + 1 * (tc_w + tc_h_gap)
    else:
        tx = tc_left + col * (tc_w + tc_h_gap)
    ty = tc_top + row * (tc_h + tc_v_gap)

    add_rect(slide7, tx, ty, tc_w, tc_h, CARD_BG, color, 2)
    txb(slide7, icon,
        tx, ty + Emu(40000), tc_w, Emu(280000),
        font_size=28, align=PP_ALIGN.CENTER)
    txb(slide7, title,
        tx, ty + Emu(300000), tc_w, Emu(220000),
        font_size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide7, desc,
        tx + Emu(40000), ty + Emu(510000),
        tc_w - Emu(80000), tc_h - Emu(560000),
        font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 하단 마무리 배너
banner_top = Emu(int(SLIDE_H - 480000))
add_rect(slide7,
    Emu(int(SLIDE_W*0.05)), banner_top,
    Emu(int(SLIDE_W*0.90)), Emu(350000),
    RGBColor(0x0A, 0x1A, 0x2A), ACCENT_CYAN, 1.5)
txb(slide7,
    "🎬  AI 영상 제작 시스템  ·  Python + Streamlit + GPT-4o Vision + FFmpeg  ·  made by s.y.Kim",
    Emu(int(SLIDE_W*0.06)), banner_top + Emu(70000),
    Emu(int(SLIDE_W*0.88)), Emu(240000),
    font_size=14, bold=True, color=ACCENT_CYAN,
    align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════
out_path = r"c:\Users\dalia\lngc\API\video_maker\AI영상제작시스템_소개.pptx"
prs.save(out_path)

import os
size = os.path.getsize(out_path)
print(f"[OK] 저장 완료: {out_path}")
print(f"     파일 크기: {size:,} bytes  ({size/1024:.1f} KB)")
